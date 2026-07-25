# build_018.py  ―  파이썬 100강 마스터 클래스 / 018강 비교 연산자
# 실행:  pip install python-pptx  →  python build_018.py
from pptx import Presentation
from pptx.util import Inches as In, Pt
from pptx.dml.color import RGBColor
from pptx.enum.shapes import MSO_SHAPE
from pptx.enum.text import PP_ALIGN, MSO_ANCHOR
from pptx.oxml.ns import qn
from lxml import etree

# ───────── Design tokens (016강과 동일) ─────────
BG, CARD, CARD2, BORDER = "0F172A", "1B2537", "162134", "334155"
YELLOW, BLUE, GREEN, ROSE = "FFD43B", "4B8BBE", "4ADE80", "FB7185"
TEXT, MUTED, FOOT, CODEBG, BAR = "E2E8F0", "94A3B8", "64748B", "111C30", "1E293B"
CMT, STR, NUM = "64748B", "FDE68A", "93C5FD"
HEAD, CODE = "Calibri", "Consolas"

def C(h): return RGBColor.from_string(h)

prs = Presentation()
prs.slide_width, prs.slide_height = In(13.333), In(7.5)
BLANK = prs.slide_layouts[6]

def new_slide():
    s = prs.slides.add_slide(BLANK)
    bg = s.shapes.add_shape(MSO_SHAPE.RECTANGLE, 0, 0, prs.slide_width, prs.slide_height)
    bg.fill.solid(); bg.fill.fore_color.rgb = C(BG)
    bg.line.fill.background(); bg.shadow.inherit = False
    return s

def txt(s, x, y, w, h, t, size=14, color=TEXT, bold=False, font=HEAD, align="l", ls=None):
    tb = s.shapes.add_textbox(In(x), In(y), In(w), In(h))
    tf = tb.text_frame; tf.word_wrap = True
    tf.margin_left = tf.margin_right = tf.margin_top = tf.margin_bottom = 0
    p = tf.paragraphs[0]
    p.alignment = {"l": PP_ALIGN.LEFT, "c": PP_ALIGN.CENTER, "r": PP_ALIGN.RIGHT}[align]
    if ls: p.line_spacing = ls
    r = p.add_run(); r.text = t
    f = r.font; f.name = font; f.size = Pt(size); f.bold = bold; f.color.rgb = C(color)
    return tb

def lines(s, x, y, w, h, items, size=14, font=HEAD, gap=7, ls=1.25):
    """items = [(text, color, bold), ...]"""
    tb = s.shapes.add_textbox(In(x), In(y), In(w), In(h))
    tf = tb.text_frame; tf.word_wrap = True
    tf.margin_left = tf.margin_right = tf.margin_top = tf.margin_bottom = 0
    for i, (t, col, bold) in enumerate(items):
        p = tf.paragraphs[0] if i == 0 else tf.add_paragraph()
        p.line_spacing = ls; p.space_after = Pt(gap)
        r = p.add_run(); r.text = t
        f = r.font; f.name = font; f.size = Pt(size); f.bold = bold; f.color.rgb = C(col)
    return tb

def card(s, x, y, w, h, fill=CARD, border=None, radius=0.05):
    sh = s.shapes.add_shape(MSO_SHAPE.ROUNDED_RECTANGLE, In(x), In(y), In(w), In(h))
    sh.adjustments[0] = radius
    sh.fill.solid(); sh.fill.fore_color.rgb = C(fill)
    if border:
        sh.line.color.rgb = C(border); sh.line.width = Pt(1)
    else:
        sh.line.fill.background()
    sh.shadow.inherit = False
    return sh

def dot(s, x, y, d, color):
    sh = s.shapes.add_shape(MSO_SHAPE.OVAL, In(x), In(y), In(d), In(d))
    sh.fill.solid(); sh.fill.fore_color.rgb = C(color)
    sh.line.fill.background(); sh.shadow.inherit = False
    return sh

def est(t, size):
    return sum((0.98 if ord(ch) > 0x1100 else 0.56) * size for ch in t) / 72

def pill(s, x, y, t, color, size=12):
    w = est(t, size) + 0.44
    sh = card(s, x, y, w, 0.36, fill=BAR, border=color, radius=0.5)
    tf = sh.text_frame; tf.word_wrap = False
    tf.margin_left = tf.margin_right = tf.margin_top = tf.margin_bottom = 0
    p = tf.paragraphs[0]; p.alignment = PP_ALIGN.CENTER
    r = p.add_run(); r.text = t
    f = r.font; f.name = CODE; f.size = Pt(size); f.bold = True; f.color.rgb = C(color)
    return w

def title_bar(s, kicker, title):
    txt(s, 0.7, 0.42, 8, 0.3, kicker, size=13, color=YELLOW, bold=True)
    txt(s, 0.7, 0.75, 11.9, 0.62, title, size=32, color=TEXT, bold=True)

def footer(s, right=""):
    txt(s, 0.7, 6.95, 8, 0.3, "파이썬 100강 마스터 클래스 · 018강 비교 연산자", size=11, color=FOOT)
    if right:
        txt(s, 9.6, 6.95, 3.0, 0.3, right, size=11, color=FOOT, align="r")

def code_panel(s, x, y, w, rows, title="python", size=13):
    """rows = [[(text, color), ...], ...]  → 반환값: 패널 높이"""
    bar, pad = 0.34, 0.16
    lh = size * 1.62 / 72
    h = bar + pad * 2 + len(rows) * lh
    card(s, x, y, w, h, fill=CODEBG, border=BORDER, radius=0.04)
    top = s.shapes.add_shape(MSO_SHAPE.RECTANGLE, In(x), In(y), In(w), In(bar))
    top.fill.solid(); top.fill.fore_color.rgb = C(BAR)
    top.line.fill.background(); top.shadow.inherit = False
    for i, c in enumerate(["FF5F56", "FFBD2E", "27C93F"]):
        dot(s, x + 0.15 + i * 0.19, y + 0.115, 0.11, c)
    txt(s, x + 0.78, y + 0.05, w - 1.0, 0.24, title, size=10, color=MUTED)
    tb = s.shapes.add_textbox(In(x + 0.22), In(y + bar + pad - 0.06), In(w - 0.4), In(h - bar - pad))
    tf = tb.text_frame; tf.word_wrap = False
    tf.margin_left = tf.margin_right = tf.margin_top = tf.margin_bottom = 0
    for i, row in enumerate(rows):
        p = tf.paragraphs[0] if i == 0 else tf.add_paragraph()
        p.line_spacing = 1.3
        for t, col in row:
            r = p.add_run(); r.text = t
            f = r.font; f.name = CODE; f.size = Pt(size); f.color.rgb = C(col)
    return h

def notes(s, t):
    s.notes_slide.notes_text_frame.text = t

# ═══════════════ 1. 타이틀 ═══════════════
s = new_slide()
deco = s.shapes.add_shape(MSO_SHAPE.OVAL, In(9.6), In(-1.4), In(5.4), In(5.4))
deco.fill.solid(); deco.fill.fore_color.rgb = C(BAR)
deco.line.fill.background(); deco.shadow.inherit = False
txt(s, 9.9, 0.85, 4.5, 1.0, "==", size=72, color=BORDER, bold=True, font=CODE, align="c")

txt(s, 0.9, 1.85, 8, 0.34, "STAGE 2 · 연산자와 제어문", size=14, color=YELLOW, bold=True)
txt(s, 0.9, 2.25, 4, 0.85, "018강", size=54, color=BLUE, bold=True)
txt(s, 0.9, 3.15, 11, 0.95, "비교 연산자", size=44, color=TEXT, bold=True)
txt(s, 0.9, 4.15, 9, 0.4, "참(True)과 거짓(False), 프로그램이 판단하는 법", size=18, color=MUTED)
px = 0.9
for t, c in [("==   !=", YELLOW), (">   <", BLUE), (">=   <=", GREEN)]:
    px += pill(s, px, 4.95, t, c, 13) + 0.25
txt(s, 0.9, 6.55, 8, 0.3, "파이썬 100강 마스터 클래스 · 왕초보", size=12, color=FOOT)
notes(s, "016강에서 계산하는 연산자를 배웠다면, 018강은 판단하는 연산자입니다. "
         "결과가 오직 True/False 두 가지뿐이라는 점이 핵심이며, 다음 강의 if문의 재료가 됩니다.")

# ═══════════════ 2. 비교 연산자란? ═══════════════
s = new_slide()
title_bar(s, "01. 개념", "비교 연산자란 무엇인가?")
card(s, 0.7, 1.70, 5.85, 1.85, fill=CARD)
txt(s, 1.0, 1.95, 5.2, 0.3, "우리의 일상", size=15, color=YELLOW, bold=True)
lines(s, 1.0, 2.42, 5.3, 1.0, [
    ("이 물건이 저 물건보다 더 비싼가?", TEXT, False),
    ("두 사람의 나이가 같은가?", TEXT, False)], size=14)
card(s, 6.75, 1.70, 5.85, 1.85, fill=CARD)
txt(s, 7.05, 1.95, 5.2, 0.3, "컴퓨터 프로그램", size=15, color=BLUE, bold=True)
lines(s, 7.05, 2.42, 5.3, 1.0, [
    ("입력한 비밀번호가 저장된 값과 일치하는가?", TEXT, False),
    ("캐릭터의 체력(HP)이 0 이하로 떨어졌는가?", TEXT, False)], size=14)

card(s, 0.7, 3.80, 11.9, 1.00, fill=CARD2, border=BORDER)
txt(s, 1.05, 4.02, 11.2, 0.5,
    "비교 연산자는 양옆의 두 값을 비교해 오직 참(True) 또는 거짓(False) 만 돌려줍니다. → 자료형은 bool",
    size=16, color=TEXT, bold=True)
code_panel(s, 0.7, 5.05, 11.9, [
    [("print", (BLUE)), ("(", TEXT), ("10 > 3", TEXT), (")", TEXT), ("      # True", GREEN)],
    [("print", (BLUE)), ("(", TEXT), ("10 == 3", TEXT), (")", TEXT), ("     # False", ROSE)],
    [("print", (BLUE)), ("(", TEXT), ("10 != 3", TEXT), (")", TEXT), ("     # True", GREEN)],
], title="compare.py")
footer(s, "02 / 10")
notes(s, "비교 연산자의 결과는 반드시 True 아니면 False 하나입니다. "
         "숫자가 나오는 산술 연산자와 결정적으로 다른 점이라고 강조해 주세요.")

# ═══════════════ 3. 6가지 표 ═══════════════
s = new_slide()
title_bar(s, "02. 종류", "비교 연산자 6가지")
txt(s, 0.7, 1.45, 8, 0.3, "a = 10,  b = 20  일 때", size=14, color=MUTED, font=CODE)

rows = [
    ("연산자", "의미", "예시", "결과"),
    ("==", "양쪽의 값이 같으면 참", "a == b", "False"),
    ("!=", "양쪽의 값이 다르면 참", "a != b", "True"),
    (">",  "왼쪽 값이 더 크면 참 (초과)", "a > b", "False"),
    ("<",  "왼쪽 값이 더 작으면 참 (미만)", "a < b", "True"),
    (">=", "왼쪽 값이 크거나 같으면 참 (이상)", "a >= b", "False"),
    ("<=", "왼쪽 값이 작거나 같으면 참 (이하)", "a <= b", "True"),
]
gf = s.shapes.add_table(7, 4, In(0.7), In(1.88), In(11.9), In(4.1))
tbl = gf.table
tbl.first_row = False; tbl.horz_banding = False
tblPr = tbl._tbl.tblPr
etree.SubElement(tblPr, qn('a:tableStyleId')).text = "{2D5ABB26-0587-4C30-8999-92F81FD0307C}"
for i, w in enumerate([1.7, 5.0, 2.7, 2.5]):
    tbl.columns[i].width = In(w)
tbl.rows[0].height = In(0.52)
for i in range(1, 7):
    tbl.rows[i].height = In(0.58)

def set_cell(cell, t, size=14, color=TEXT, bold=False, font=HEAD, fill=CARD, align="c"):
    cell.fill.solid(); cell.fill.fore_color.rgb = C(fill)
    cell.vertical_anchor = MSO_ANCHOR.MIDDLE
    cell.margin_left = cell.margin_right = In(0.14)
    cell.margin_top = cell.margin_bottom = In(0.04)
    p = cell.text_frame.paragraphs[0]
    p.alignment = {"l": PP_ALIGN.LEFT, "c": PP_ALIGN.CENTER}[align]
    r = p.add_run(); r.text = t
    f = r.font; f.name = font; f.size = Pt(size); f.bold = bold; f.color.rgb = C(color)

for ri, row in enumerate(rows):
    head = (ri == 0)
    bgc = BAR if head else (CARD if ri % 2 else CARD2)
    for ci, val in enumerate(row):
        if head:
            set_cell(tbl.cell(ri, ci), val, size=14, color=YELLOW, bold=True, fill=BAR)
        elif ci == 0:
            set_cell(tbl.cell(ri, ci), val, size=17, color=YELLOW, bold=True, font=CODE, fill=bgc)
        elif ci == 1:
            set_cell(tbl.cell(ri, ci), val, size=14, color=TEXT, fill=bgc, align="l")
        elif ci == 2:
            set_cell(tbl.cell(ri, ci), val, size=14, color=MUTED, font=CODE, fill=bgc)
        else:
            set_cell(tbl.cell(ri, ci), val, size=14, color=(GREEN if val == "True" else ROSE),
                     bold=True, font=CODE, fill=bgc)

txt(s, 0.7, 6.20, 11.9, 0.3,
    "수학의  ≥  ≤  ≠  기호는 키보드에 없습니다. 그래서 파이썬은  >=   <=   !=  로 대신 씁니다.",
    size=13, color=MUTED)
footer(s, "03 / 10")
notes(s, "표를 한 줄씩 읽으며 a=10, b=20을 대입해 결과를 함께 확인합니다. "
         "부등호가 먼저, 등호(=)가 뒤에 온다는 순서(>=, <=)를 꼭 짚어 주세요.")

# ═══════════════ 4. > vs >= 경계값 ═══════════════
s = new_slide()
title_bar(s, "03. 헷갈리는 경계", "> 와 >= 는 무엇이 다른가?")
card(s, 0.7, 1.70, 5.85, 3.50, fill=CARD)
txt(s, 1.0, 1.95, 5.2, 0.35, ">   초과", size=20, color=YELLOW, bold=True, font=CODE)
txt(s, 1.0, 2.45, 5.3, 0.35, "기준값 자기 자신은 포함하지 않는다", size=14, color=MUTED)
code_panel(s, 1.0, 2.95, 5.25, [
    [("x = ", TEXT), ("10", NUM)],
    [("print", BLUE), ("(x > ", TEXT), ("10", NUM), (")", TEXT)],
    [("# False", ROSE)],
], title="gt.py", size=12)

card(s, 6.75, 1.70, 5.85, 3.50, fill=CARD)
txt(s, 7.05, 1.95, 5.2, 0.35, ">=   이상", size=20, color=GREEN, bold=True, font=CODE)
txt(s, 7.05, 2.45, 5.3, 0.35, "기준값 자기 자신을 포함한다", size=14, color=MUTED)
code_panel(s, 7.05, 2.95, 5.25, [
    [("x = ", TEXT), ("10", NUM)],
    [("print", BLUE), ("(x >= ", TEXT), ("10", NUM), (")", TEXT)],
    [("# True", GREEN)],
], title="ge.py", size=12)

card(s, 0.7, 5.45, 11.9, 1.10, fill=CARD2, border=BORDER)
dot(s, 1.02, 5.86, 0.24, YELLOW)
txt(s, 1.42, 5.60, 11.0, 0.85,
    "성인 판별을  age >= 19  로 쓰면 19세가 포함되고,  age > 19  로 쓰면 19세가 제외됩니다.\n"
    "경계값 하나가 서비스 정책을 통째로 바꿉니다. '이상/이하'인지 '초과/미만'인지 항상 확인하세요.",
    size=13, color=TEXT, ls=1.3)
footer(s, "04 / 10")
notes(s, "초보자가 논리 버그를 가장 많이 만드는 지점입니다. "
         "'19세도 성인인가?'처럼 경계값을 말로 먼저 확인하는 습관을 강조하세요.")

# ═══════════════ 5. = 와 == ═══════════════
s = new_slide()
title_bar(s, "04. 초보자 최대 실수", "=  와  ==  는 완전히 다르다")
card(s, 0.7, 1.70, 5.85, 3.15, fill=CARD)
txt(s, 1.0, 1.95, 5.2, 0.4, "=     대입 연산자", size=19, color=BLUE, bold=True, font=CODE)
txt(s, 1.0, 2.45, 5.3, 0.35, "\"오른쪽 값을 왼쪽 변수에 넣어라\"", size=14, color=MUTED)
code_panel(s, 1.0, 2.95, 5.25, [
    [("x = ", TEXT), ("10", NUM), ("        # 저장(명령)", CMT)],
], title="assign.py", size=12)

card(s, 6.75, 1.70, 5.85, 3.15, fill=CARD)
txt(s, 7.05, 1.95, 5.2, 0.4, "==    비교 연산자", size=19, color=YELLOW, bold=True, font=CODE)
txt(s, 7.05, 2.45, 5.3, 0.35, "\"양쪽 값이 같은지 물어보라\"", size=14, color=MUTED)
code_panel(s, 7.05, 2.95, 5.25, [
    [("x == ", TEXT), ("10", NUM), ("       # 질문(True)", CMT)],
], title="compare.py", size=12)

card(s, 0.7, 5.10, 11.9, 1.45, fill=CARD2, border=ROSE)
dot(s, 1.02, 5.44, 0.26, ROSE)
txt(s, 1.45, 5.28, 4.4, 0.35, "= 하나만 쓰면?", size=15, color=ROSE, bold=True)
txt(s, 1.45, 5.68, 5.5, 0.7,
    "파이썬이 '대입 명령'으로 읽어\n문법 오류(SyntaxError)를 냅니다.", size=13, color=TEXT, ls=1.3)
code_panel(s, 7.05, 5.25, 5.25, [
    [("if", BLUE), (" x = ", TEXT), ("10", NUM), (":", TEXT), ("   # SyntaxError", ROSE)],
], title="error.py", size=12)
footer(s, "05 / 10")
notes(s, "= 는 명령문, == 는 질문입니다. 이 한 문장으로 정리해 주세요. "
         "실제로 코드를 실행해 SyntaxError 메시지를 눈으로 보여주면 기억에 오래 남습니다.")

# ═══════════════ 6. 결과는 bool ═══════════════
s = new_slide()
title_bar(s, "05. 결과의 정체", "비교의 결과는 언제나 bool 이다")
code_panel(s, 0.7, 1.70, 11.9, [
    [("result = ", TEXT), ("10", NUM), (" > ", TEXT), ("3", NUM)],
    [("print", BLUE), ("(result)", TEXT), ("            # True", GREEN)],
    [("print", BLUE), ("(", TEXT), ("type", BLUE), ("(result))", TEXT), ("      # <class 'bool'>", CMT)],
], title="bool.py", size=14)

card(s, 0.7, 3.95, 5.85, 2.10, fill=CARD)
txt(s, 1.0, 4.20, 5.2, 0.32, "012강 type() 과 연결", size=15, color=BLUE, bold=True)
txt(s, 1.0, 4.62, 5.3, 1.2,
    "어떤 값을 비교하든 결과의 타입은 늘 bool 하나입니다.\n"
    "숫자를 비교해도, 문자열을 비교해도 마찬가지입니다.", size=13, color=TEXT, ls=1.3)

card(s, 6.75, 3.95, 5.85, 2.10, fill=CARD)
txt(s, 7.05, 4.20, 5.2, 0.32, "대소문자에 주의", size=15, color=YELLOW, bold=True)
txt(s, 7.05, 4.62, 5.3, 1.2,
    "True / False 는 첫 글자가 반드시 대문자입니다.\n"
    "true, TRUE 라고 쓰면 NameError 가 납니다.", size=13, color=TEXT, ls=1.3)
footer(s, "06 / 10")
notes(s, "12강에서 배운 type() 함수를 다시 꺼내 결과가 bool임을 눈으로 확인시킵니다. "
         "True/False 대소문자 실수는 실제로 매우 흔합니다.")

# ═══════════════ 7. 문자열 비교 ═══════════════
s = new_slide()
title_bar(s, "06. 확장", "숫자만? 문자열도 비교할 수 있다")
code_panel(s, 0.7, 1.70, 11.9, [
    [("print", BLUE), ("(", TEXT), ("\"apple\"", STR), (" == ", TEXT), ("\"Apple\"", STR), (")", TEXT),
     ("    # False  ← 대소문자를 구분한다", CMT)],
    [("print", BLUE), ("(", TEXT), ("\"apple\"", STR), (" == ", TEXT), ("\"apple\"", STR), (")", TEXT),
     ("    # True", GREEN)],
    [("print", BLUE), ("(", TEXT), ("\"apple\"", STR), (" < ", TEXT), ("\"banana\"", STR), (")", TEXT),
     ("   # True   ← 사전(알파벳) 순서로 비교", CMT)],
], title="string_compare.py", size=13)

card(s, 0.7, 4.05, 11.9, 2.35, fill=CARD2, border=BORDER)
dot(s, 1.02, 4.45, 0.24, ROSE)
txt(s, 1.42, 4.28, 11.0, 0.35, "타입이 다르면?", size=16, color=ROSE, bold=True)
lines(s, 1.42, 4.78, 11.0, 1.4, [
    ("print(10 == \"10\")   →   False       숫자 10과 문자열 \"10\"은 서로 다른 값입니다.", TEXT, False),
    ("print(10 > \"5\")     →   TypeError   크기 비교는 아예 불가능합니다.", TEXT, False),
    ("사용자 입력값(input)은 항상 문자열이므로, 숫자와 비교하려면 int()로 바꿔야 합니다.", MUTED, False),
], size=13, font=HEAD, gap=8)
footer(s, "07 / 10")
notes(s, "== 로 같은지는 물어볼 수 있지만(False), 다른 타입끼리 크기 비교는 TypeError입니다. "
         "input()과 형변환 이야기의 복선입니다.")

# ═══════════════ 8. 실전 활용 ═══════════════
s = new_slide()
title_bar(s, "07. 실전 활용", "실제로 이럴 때 씁니다")
txt(s, 0.7, 1.62, 5.85, 0.32, "① 로그인 비밀번호 검증", size=15, color=YELLOW, bold=True)
code_panel(s, 0.7, 2.05, 5.85, [
    [("saved_pw = ", TEXT), ("\"python123\"", STR)],
    [("input_pw = ", TEXT), ("\"python123\"", STR)],
    [("", TEXT)],
    [("print", BLUE), ("(saved_pw == input_pw)", TEXT)],
    [("# True → 로그인 성공", GREEN)],
], title="login.py", size=12)

txt(s, 6.75, 1.62, 5.85, 0.32, "② 게임 체력(HP) 판정", size=15, color=BLUE, bold=True)
code_panel(s, 6.75, 2.05, 5.85, [
    [("hp = ", TEXT), ("0", NUM)],
    [("", TEXT)],
    [("", TEXT)],
    [("print", BLUE), ("(hp <= ", TEXT), ("0", NUM), (")", TEXT)],
    [("# True → 게임 오버", ROSE)],
], title="game.py", size=12)

card(s, 0.7, 5.10, 11.9, 1.30, fill=CARD2, border=BORDER)
dot(s, 1.02, 5.62, 0.24, GREEN)
txt(s, 1.42, 5.32, 11.0, 0.9,
    "지금은 True / False 를 화면에 출력만 했습니다.\n"
    "곧 배울 if 문과 만나는 순간, 이 값은 프로그램의 흐름을 바꾸는 '스위치'가 됩니다.",
    size=14, color=TEXT, ls=1.35)
footer(s, "08 / 10")
notes(s, "비교 연산자가 왜 필요한지를 실제 사례로 못 박는 슬라이드입니다. "
         "if문 예고로 자연스럽게 넘어갑니다.")

# ═══════════════ 9. 3대 함정 ═══════════════
s = new_slide()
title_bar(s, "08. 이것만은 조심", "비교 연산자 3대 함정")
traps = [
    (0.70, "1", "!= 를  =! 로 쓰기", ["느낌표가 반드시 앞!", "a != b   (O)", "a =! b   (X)"], ROSE),
    (4.75, "2", "기호 사이 띄어쓰기", [">= 는 붙여 씁니다.", "a >= b   (O)", "a > = b  (X) SyntaxError"], YELLOW),
    (8.80, "3", "실수(float) 비교", ["0.1 + 0.2 == 0.3", "→ False (부동소수점 오차)", "9강 내용 복습!"], BLUE),
]
for x, num, head, body, col in traps:
    card(s, x, 1.70, 3.83, 3.55, fill=CARD)
    dot(s, x + 0.28, 1.98, 0.34, col)
    txt(s, x + 0.28, 2.02, 0.34, 0.28, num, size=13, color=BG, bold=True, align="c")
    txt(s, x + 0.28, 2.52, 3.3, 0.35, head, size=15, color=col, bold=True)
    lines(s, x + 0.28, 3.05, 3.3, 1.9, [(b, (MUTED if i == 0 else TEXT), False) for i, b in enumerate(body)],
          size=12, font=CODE, gap=9)

card(s, 0.7, 5.50, 11.9, 1.05, fill=CARD2, border=BORDER)
txt(s, 1.05, 5.72, 11.2, 0.6,
    "실수를 비교할 때는  a == b  대신  abs(a - b) < 0.000001  처럼 오차를 감안해야 안전합니다.",
    size=14, color=TEXT)
footer(s, "09 / 10")
notes(s, "세 함정 모두 실제 수업에서 반복적으로 나오는 질문입니다. "
         "특히 0.1 + 0.2 == 0.3 은 직접 실행해 보여주면 반응이 큽니다.")

# ═══════════════ 10. 핵심 정리 ═══════════════
s = new_slide()
title_bar(s, "핵심 정리", "018강 요약")
summary = [
    "비교 연산자는 6가지 —  ==   !=   >   <   >=   <=",
    "결과는 언제나 True 또는 False 단 두 가지 (자료형: bool)",
    "=  는 대입(명령),  ==  는 비교(질문) — 가장 흔한 실수",
    ">(초과) 와 >=(이상) 의 경계값 차이에 항상 주의",
]
y = 1.80
for i, t in enumerate(summary):
    card(s, 0.7, y, 11.9, 0.82, fill=(CARD if i % 2 == 0 else CARD2))
    dot(s, 1.05, y + 0.33, 0.16, YELLOW)
    txt(s, 1.45, y + 0.22, 11.0, 0.4, t, size=15, color=TEXT)
    y += 0.95

card(s, 0.7, 5.75, 11.9, 1.00, fill=BAR, border=BLUE)
txt(s, 1.05, 5.92, 2.2, 0.35, "019강 예고", size=14, color=BLUE, bold=True)
txt(s, 3.15, 5.92, 9.2, 0.35,
    "논리 연산자 (and, or, not) — 여러 개의 조건을 하나로 묶는 법", size=15, color=TEXT, bold=True)
notes(s, "요약 4줄을 학습자가 스스로 말해보게 하면 좋습니다. "
         "019강에서는 조건을 여러 개 결합하는 and / or / not 을 다룹니다.")

prs.save("018강_비교연산자.pptx")
print("생성 완료 → 018강_비교연산자.pptx")